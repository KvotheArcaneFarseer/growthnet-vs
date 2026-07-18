#!/usr/bin/env python3
"""
Build a PowerPoint deck for the GrowthNet lab meeting.

Assembles:
  - Title slide
  - Pipeline overview slide
  - One case slide per seed variant (2D early + late QC panels)
  - 3D napari visualization slide

Inputs:  existing PNG outputs (no pipeline re-run needed)
Output:  lab_meeting_exports/GrowthNet_LabMeeting_2026-04-25.pptx

Usage:
    python3 scripts/build_lab_meeting_deck.py

Prerequisites:
    pip install python-pptx
    (optional) python3 scripts/export_3d_frames.py  # for 3D slide
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: python3 -m pip install python-pptx")
    sys.exit(1)

REPO_ROOT = Path(__file__).resolve().parents[1]
ANIM_DIR = REPO_ROOT / "animation_outputs"
OUT_DIR = REPO_ROOT / "lab_meeting_exports"
FRAMES_DIR = OUT_DIR / "3d_frames"

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5

# Case definitions — (label, early_png, late_png)
CASES: list[tuple[str, Path, Path]] = [
    (
        "Sample A (auto seed α)",
        REPO_ROOT / "tmp_seed_validation/auto_seed_a/qc_embedding.png",
        REPO_ROOT / "tmp_seed_validation/auto_seed_a/qc_embedding_late.png",
    ),
    (
        "Sample B (auto seed β)",
        REPO_ROOT / "tmp_seed_validation/auto_seed_b/qc_embedding.png",
        REPO_ROOT / "tmp_seed_validation/auto_seed_b/qc_embedding_late.png",
    ),
    (
        "Sample C (alternative seed)",
        REPO_ROOT / "tmp_seed_validation/diff_seed/qc_embedding.png",
        REPO_ROOT / "tmp_seed_validation/diff_seed/qc_embedding_late.png",
    ),
    (
        "Sample D (fixed seed, run 1)",
        REPO_ROOT / "tmp_seed_validation/same_seed_a/qc_embedding.png",
        REPO_ROOT / "tmp_seed_validation/same_seed_a/qc_embedding_late.png",
    ),
    (
        "Sample E (direct single-case run)",
        REPO_ROOT / "embedding_outputs/qc_embedding.png",
        REPO_ROOT / "embedding_outputs/qc_embedding_late.png",
    ),
]

BLACK = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x8B, 0xFF)


def _in(val: float) -> Emu:
    return Inches(val)


def _add_text(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    txb = slide.shapes.add_textbox(_in(left), _in(top), _in(width), _in(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _fill_slide_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_title_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _fill_slide_bg(slide, RGBColor(0x0D, 0x1B, 0x2A))

    _add_text(slide, "GrowthNet", 0.6, 1.2, 12.0, 1.4,
              font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, "Synthetic Vestibular Schwannoma Embedding",
              0.6, 2.6, 12.0, 0.8,
              font_size=28, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    _add_text(slide, "Current Pipeline Results — April 25, 2026",
              0.6, 3.5, 12.0, 0.6,
              font_size=18, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
    _add_text(slide, "Prepared for: Ethan · Sofia · George",
              0.6, 4.3, 12.0, 0.5,
              font_size=16, color=RGBColor(0x66, 0x88, 0xAA), align=PP_ALIGN.CENTER)


def make_overview_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_slide_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))

    _add_text(slide, "What We're Showing", 0.5, 0.3, 12.0, 0.7,
              font_size=32, bold=True, color=RGBColor(0x0D, 0x1B, 0x2A))

    body = (
        "Each panel shows a 3×3 QC grid:\n"
        "  Rows:     Axial  /  Coronal  /  Sagittal slice\n"
        "  Columns:  Original MRI  |  Real tumour segmentation (red)  |  Embedded synthetic tumour (cyan)\n\n"
        "Pipeline:  Lollipop geometry → principal-axis orientation → rigid placement into real MRI\n"
        "Patient:   Case 147_0_0  (single T2 MRI with labelled VS tumour)\n"
        "Samples:   5 independent synthetic tumour generations (varied seeds / placement)\n"
        "Timepoints per sample:  t=0 (smallest) and t=4 (elongated late stage)\n\n"
        "3D visualisation:  napari renders of synthetic tumour mask (see last slide)"
    )
    _add_text(slide, body, 0.5, 1.2, 12.3, 5.5,
              font_size=17, color=RGBColor(0x20, 0x20, 0x30))


def make_case_slide(prs: Presentation, label: str, early_png: Path, late_png: Path) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_slide_bg(slide, RGBColor(0xF8, 0xF9, 0xFF))

    _add_text(slide, label, 0.3, 0.1, 12.5, 0.55,
              font_size=22, bold=True, color=RGBColor(0x0D, 0x1B, 0x2A))

    margin = 0.25
    panel_w = (SLIDE_W_IN - 3 * margin) / 2
    panel_h = SLIDE_H_IN - 1.1

    for i, (png, sublabel) in enumerate([(early_png, "t=0  (early)"), (late_png, "t=4  (late / elongated)")]):
        left = margin + i * (panel_w + margin)
        top = 0.85
        if png.exists():
            slide.shapes.add_picture(str(png), _in(left), _in(top), _in(panel_w), _in(panel_h))
        else:
            _add_text(slide, f"[missing: {png.name}]", left, top + panel_h / 2, panel_w, 0.5,
                      font_size=12, color=RGBColor(0xAA, 0x00, 0x00))
        _add_text(slide, sublabel, left, top + panel_h + 0.03, panel_w, 0.3,
                  font_size=13, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)


def make_3d_slide(prs: Presentation) -> bool:
    angle1 = FRAMES_DIR / "3d_angle1.png"
    angle2 = FRAMES_DIR / "3d_angle2.png"

    has_angle1 = angle1.exists() and angle1.stat().st_size > 0
    has_angle2 = angle2.exists() and angle2.stat().st_size > 0

    if not has_angle1 and not has_angle2:
        return False

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _fill_slide_bg(slide, RGBColor(0x0D, 0x1B, 0x2A))

    _add_text(slide, "3D Tumour Visualisation  (napari)", 0.3, 0.1, 12.5, 0.65,
              font_size=26, bold=True, color=WHITE)
    _add_text(slide, "Synthetic lollipop-geometry VS mask rendered in napari",
              0.3, 0.72, 12.5, 0.35, font_size=15, color=RGBColor(0x88, 0xAA, 0xCC))

    top = 1.15
    panel_h = SLIDE_H_IN - top - 0.25

    if has_angle1 and has_angle2:
        margin = 0.25
        panel_w = (SLIDE_W_IN - 3 * margin) / 2
        for i, (png, lbl) in enumerate([(angle1, "Angle 1"), (angle2, "Angle 2")]):
            left = margin + i * (panel_w + margin)
            slide.shapes.add_picture(str(png), _in(left), _in(top), _in(panel_w), _in(panel_h))
            _add_text(slide, lbl, left, top + panel_h + 0.02, panel_w, 0.25,
                      font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
    else:
        png = angle1 if has_angle1 else angle2
        lbl = "Angle 1" if has_angle1 else "Angle 2"
        panel_w = SLIDE_W_IN * 0.7
        left = (SLIDE_W_IN - panel_w) / 2
        slide.shapes.add_picture(str(png), _in(left), _in(top), _in(panel_w), _in(panel_h))
        _add_text(slide, lbl, left, top + panel_h + 0.02, panel_w, 0.25,
                  font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

    return True


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # Try to extract 3D frames first if not present
    angle1 = FRAMES_DIR / "3d_angle1.png"
    angle2 = FRAMES_DIR / "3d_angle2.png"
    if not (angle1.exists() and angle2.exists()):
        print("3D frames not found — running export_3d_frames.py...")
        extractor = Path(__file__).parent / "export_3d_frames.py"
        if extractor.exists():
            subprocess.run([sys.executable, str(extractor)], check=False)
        else:
            print("  export_3d_frames.py not found — 3D slide will be skipped if frames missing")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    print("Building slide deck...")

    make_title_slide(prs)
    print("  [1] Title slide")

    make_overview_slide(prs)
    print("  [2] Overview slide")

    n_cases = 0
    for label, early, late in CASES:
        if early.exists() or late.exists():
            make_case_slide(prs, label, early, late)
            n_cases += 1
            print(f"  [{2 + n_cases}] {label}")
        else:
            print(f"  SKIP (missing PNGs): {label}")

    has_3d = make_3d_slide(prs)
    if has_3d:
        print(f"  [{3 + n_cases}] 3D napari slide")
    else:
        print("  SKIP: 3D slide (no frames found — run export_3d_frames.py first)")

    out_path = OUT_DIR / "GrowthNet_LabMeeting_2026-04-25.pptx"
    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    print(f"Slides: 1 title + 1 overview + {n_cases} cases + {'1 3D' if has_3d else '0 3D'} = {2 + n_cases + (1 if has_3d else 0)} total")


if __name__ == "__main__":
    main()
